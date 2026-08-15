import logfire


def _parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Include table cell text (some documents store content in tables)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def parse_office(file_path: str):
    """
    Parses Office documents (.docx via python-docx, .pptx via python-pptx).

    Note: uses these libraries directly instead of unstructured's auto-partition,
    which can hang indefinitely on some .pptx files.
    """
    with logfire.span("📄 Office Document Parsing", filename=file_path):
        try:
            ext = file_path.lower().rsplit(".", 1)[-1]
            if ext == "docx":
                full_text = _parse_docx(file_path)
            elif ext == "pptx":
                full_text = _parse_pptx(file_path)
            else:
                logfire.warning(f"Unsupported office extension for {file_path}")
                return ""

            if not full_text.strip():
                logfire.warning(f"⚠️ No text extracted from {file_path}")
            else:
                logfire.info(f"✅ Successfully parsed {len(full_text)} characters")

            return full_text
        except Exception as e:
            logfire.error(f"❌ Office Parse Failed: {e}")
            raise e
